"""Runtime document management + re-indexing for the deployed chatbot.

Lets Streamlit users add/delete docs and rebuild the FAISS index without redeploying.
"""
from __future__ import annotations

import json
import os
import re
import shutil
from pathlib import Path
from typing import Iterable

ALLOWED_EXTENSIONS = {".txt", ".md", ".csv", ".pdf", ".json", ".xlsx", ".docx", ".pptx"}


def sanitize_filename(name: str) -> str:
    """Strip path components + restrict charset; raise on unsupported extension."""
    stem = Path(name).stem
    suffix = Path(name).suffix.lower()
    safe_stem = re.sub(r"[^A-Za-z0-9가-힣._-]+", "-", stem).strip(".-_")[:80]
    if not safe_stem:
        safe_stem = "file"
    if suffix not in ALLOWED_EXTENSIONS:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {suffix}. 허용: {sorted(ALLOWED_EXTENSIONS)}")
    return f"{safe_stem}{suffix}"


def list_documents(docs_dir: Path) -> list[dict]:
    """List documents currently indexed by this chatbot."""
    docs_dir.mkdir(parents=True, exist_ok=True)
    items = []
    for p in sorted(docs_dir.iterdir()):
        if p.is_file() and p.suffix.lower() in ALLOWED_EXTENSIONS:
            try:
                size_kb = max(1, p.stat().st_size // 1024)
            except OSError:
                size_kb = 0
            items.append({"name": p.name, "size_kb": size_kb, "path": str(p)})
    return items


def save_uploaded_file(file_obj, filename: str, docs_dir: Path) -> Path:
    """Save a Streamlit/FastAPI-style UploadedFile object to docs_dir."""
    safe_name = sanitize_filename(filename)
    docs_dir.mkdir(parents=True, exist_ok=True)
    target = docs_dir / safe_name
    data = file_obj.read() if hasattr(file_obj, "read") else file_obj
    if isinstance(data, str):
        data = data.encode("utf-8")
    target.write_bytes(data)
    return target


def delete_document(name: str, docs_dir: Path) -> bool:
    """Delete a single doc by filename. Returns True if removed."""
    safe = Path(name).name  # strip any path parts
    p = docs_dir / safe
    if p.exists() and p.is_file():
        p.unlink()
        return True
    return False


def _read_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        out = []
        for i, page in enumerate(reader.pages):
            try:
                out.append(f"[page {i+1}]\n" + (page.extract_text() or ""))
            except Exception:
                pass
        return "\n\n".join(out)
    if suffix == ".docx":
        import docx2txt
        return docx2txt.process(str(path)) or ""
    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        out = []
        for i, slide in enumerate(prs.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            lines.append(text)
            if lines:
                out.append(f"[slide {i}]\n" + "\n".join(lines))
        return "\n\n".join(out)
    if suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(filename=str(path), read_only=True, data_only=True)
        sheets_out = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v).strip() for v in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                sheets_out.append(f"[sheet: {ws.title}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets_out)
    return ""


def rebuild_index(
    docs_dir: Path,
    index_dir: Path,
    chunk_size: int = 900,
    chunk_overlap: int = 120,
    embedding_model: str = "text-embedding-3-small",
    openai_api_key: str | None = None,
) -> dict:
    """Re-read all docs in docs_dir → chunk → embed → save FAISS + chunks.jsonl.
    Returns a status dict.
    """
    docs_dir.mkdir(parents=True, exist_ok=True)
    index_dir.mkdir(parents=True, exist_ok=True)

    files = [p for p in docs_dir.iterdir() if p.is_file() and p.suffix.lower() in ALLOWED_EXTENSIONS]
    if not files:
        return {"ok": False, "error": "문서가 비어있습니다. 파일을 추가하세요."}

    # Load + concatenate
    from langchain_core.documents import Document
    raw_docs = []
    for p in files:
        try:
            text = _read_text(p)
        except Exception as exc:
            raw_docs.append(Document(page_content=f"[읽기 실패: {exc}]", metadata={"source": p.name, "error": True}))
            continue
        if text.strip():
            raw_docs.append(Document(page_content=text, metadata={"source": p.name}))

    # Chunk
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_documents(raw_docs)
    if not chunks:
        return {"ok": False, "error": "청크가 생성되지 않았습니다 (모든 문서가 비어있음)."}

    # Write chunks.jsonl (used by BM25 fallback)
    chunks_path = index_dir / "chunks.jsonl"
    with chunks_path.open("w", encoding="utf-8") as f:
        for i, doc in enumerate(chunks):
            f.write(json.dumps({
                "id": i,
                "text": doc.page_content,
                "metadata": doc.metadata,
            }, ensure_ascii=False) + "\n")

    # Build FAISS
    api_key = openai_api_key or os.getenv("OPENAI_API_KEY")
    if not api_key:
        # Still useful: BM25-only retrieval works without FAISS.
        # Remove stale faiss dir so retrieval falls back cleanly.
        faiss_dir = index_dir / "faiss"
        if faiss_dir.exists():
            shutil.rmtree(faiss_dir)
        (index_dir / "index_metadata.json").write_text(
            json.dumps({"chunks": len(chunks), "documents": len(files), "embedding_provider": "none"}, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return {"ok": True, "documents": len(files), "chunks": len(chunks), "faiss_built": False,
                "message": "OPENAI_API_KEY가 없어 FAISS는 건너뛰고 BM25 인덱스만 만들었습니다."}

    try:
        from langchain_community.vectorstores import FAISS
        from langchain_openai import OpenAIEmbeddings
        embeddings = OpenAIEmbeddings(
            model=embedding_model,
            api_key=api_key,
            chunk_size=256,        # 한 API 요청에 묶을 텍스트 수 (배치 임베딩)
            max_retries=5,
            request_timeout=30.0,
        )
        db = FAISS.from_documents(chunks, embeddings)
        faiss_dir = index_dir / "faiss"
        if faiss_dir.exists():
            shutil.rmtree(faiss_dir)
        db.save_local(str(faiss_dir))
        meta = {
            "chunks": len(chunks),
            "documents": len(files),
            "embedding_provider": "openai",
            "embedding_model": embedding_model,
        }
        (index_dir / "index_metadata.json").write_text(
            json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8",
        )
        return {"ok": True, "documents": len(files), "chunks": len(chunks), "faiss_built": True,
                "message": f"✅ {len(files)}개 문서 → {len(chunks)}개 청크 (FAISS 빌드 완료)"}
    except Exception as exc:
        return {"ok": False, "error": f"FAISS 빌드 실패: {exc}"}
